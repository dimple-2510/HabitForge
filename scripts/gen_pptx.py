#!/usr/bin/env python3
"""Generate HabitForge presentation PPT for PEP 4 (10 marks)"""

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
import os

prs = Presentation()
prs.slide_width = Inches(13.333)
prs.slide_height = Inches(7.5)

# Color scheme
BLUE = RGBColor(0x3B, 0x82, 0xF6)
DARK = RGBColor(0x1F, 0x29, 0x37)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
LIGHT_GRAY = RGBColor(0xF3, 0xF4, 0xF6)
ACCENT_GREEN = RGBColor(0x10, 0xB9, 0x81)
ACCENT_PURPLE = RGBColor(0x8B, 0x5C, 0xF6)
ACCENT_ORANGE = RGBColor(0xF5, 0x9E, 0x0B)
TEXT_DARK = RGBColor(0x11, 0x18, 0x27)
TEXT_GRAY = RGBColor(0x6B, 0x72, 0x80)

def add_bg(slide, color):
    bg = slide.background
    fill = bg.fill
    fill.solid()
    fill.fore_color.rgb = color

def add_shape(slide, left, top, width, height, color, alpha=None):
    shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, left, top, width, height)
    shape.fill.solid()
    shape.fill.fore_color.rgb = color
    shape.line.fill.background()
    return shape

def add_textbox(slide, left, top, width, height, text, font_size=18, bold=False, color=TEXT_DARK, alignment=PP_ALIGN.LEFT, font_name="Calibri"):
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = Pt(font_size)
    p.font.bold = bold
    p.font.color.rgb = color
    p.font.name = font_name
    p.alignment = alignment
    return txBox

def add_bullet_list(slide, left, top, width, height, items, font_size=16, color=TEXT_DARK):
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = True
    for i, item in enumerate(items):
        if i == 0:
            p = tf.paragraphs[0]
        else:
            p = tf.add_paragraph()
        p.text = item
        p.font.size = Pt(font_size)
        p.font.color.rgb = color
        p.font.name = "Calibri"
        p.space_after = Pt(8)
    return txBox

# ============================================================
# SLIDE 1: Title Slide
# ============================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])  # Blank
add_bg(slide, DARK)

# Accent bar at top
add_shape(slide, Inches(0), Inches(0), Inches(13.333), Inches(0.08), BLUE)

# Title
add_textbox(slide, Inches(1), Inches(1.5), Inches(11), Inches(1.2),
    "HabitForge", font_size=54, bold=True, color=WHITE, alignment=PP_ALIGN.CENTER)

add_textbox(slide, Inches(1), Inches(2.7), Inches(11), Inches(0.8),
    "A Full-Stack Habit Tracking Web Application", font_size=28, color=BLUE, alignment=PP_ALIGN.CENTER)

# Decorative line
add_shape(slide, Inches(5.5), Inches(3.6), Inches(2.3), Inches(0.04), ACCENT_GREEN)

# Subtitle
add_textbox(slide, Inches(1), Inches(3.9), Inches(11), Inches(0.6),
    "MCA 4th Semester Major Project | Subject Code: 23ONMCR-753 | 12 Credits",
    font_size=20, color=TEXT_GRAY, alignment=PP_ALIGN.CENTER)

# Tech stack badges
tech_items = ["Next.js 14", "Supabase", "PostgreSQL", "Tailwind CSS", "TypeScript", "Vercel"]
badge_y = Inches(4.8)
badge_x = Inches(2.5)
for tech in tech_items:
    shape = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, badge_x, badge_y, Inches(1.2), Inches(0.4))
    shape.fill.solid()
    shape.fill.fore_color.rgb = RGBColor(0x2D, 0x37, 0x48)
    shape.line.color.rgb = BLUE
    shape.line.width = Pt(1)
    tf = shape.text_frame
    tf.word_wrap = False
    p = tf.paragraphs[0]
    p.text = tech
    p.font.size = Pt(12)
    p.font.color.rgb = WHITE
    p.font.name = "Calibri"
    p.alignment = PP_ALIGN.CENTER
    tf.paragraphs[0].space_before = Pt(0)
    tf.paragraphs[0].space_after = Pt(0)
    badge_x += Inches(1.4)

# Footer
add_textbox(slide, Inches(1), Inches(6.8), Inches(11), Inches(0.4),
    "Centre for Distance & Online Education, Chandigarh University",
    font_size=14, color=TEXT_GRAY, alignment=PP_ALIGN.CENTER)

# ============================================================
# SLIDE 2: Agenda
# ============================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide, WHITE)
add_shape(slide, Inches(0), Inches(0), Inches(13.333), Inches(0.08), BLUE)

add_textbox(slide, Inches(0.8), Inches(0.5), Inches(11), Inches(0.8),
    "Agenda", font_size=36, bold=True, color=DARK)

add_shape(slide, Inches(0.8), Inches(1.3), Inches(0.15), Inches(0.3), BLUE)

agenda_items = [
    ("01", "Introduction & Problem Statement", "Why HabitForge?"),
    ("02", "Objectives & Scope", "What we aim to achieve"),
    ("03", "Technology Stack", "Tools and frameworks used"),
    ("04", "System Architecture", "High-level design overview"),
    ("05", "Key Features", "Core functionality walkthrough"),
    ("06", "Database Design", "Schema and relationships"),
    ("07", "Implementation Highlights", "Code snippets and patterns"),
    ("08", "Testing & Results", "Quality assurance approach"),
    ("09", "Live Demo", "Application walkthrough"),
    ("10", "Conclusion & Future Scope", "Summary and roadmap"),
]

y = Inches(1.8)
for num, title, desc in agenda_items:
    add_textbox(slide, Inches(0.8), y, Inches(0.8), Inches(0.5), num, font_size=18, bold=True, color=BLUE)
    add_textbox(slide, Inches(1.6), y, Inches(5), Inches(0.4), title, font_size=16, bold=True, color=TEXT_DARK)
    add_textbox(slide, Inches(1.6), y + Inches(0.35), Inches(5), Inches(0.3), desc, font_size=12, color=TEXT_GRAY)
    y += Inches(0.75)

# ============================================================
# SLIDE 3: Introduction & Problem Statement
# ============================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide, WHITE)
add_shape(slide, Inches(0), Inches(0), Inches(13.333), Inches(0.08), BLUE)

add_textbox(slide, Inches(0.8), Inches(0.5), Inches(11), Inches(0.8),
    "Introduction & Problem Statement", font_size=32, bold=True, color=DARK)

# Problem box
prob = add_shape(slide, Inches(0.8), Inches(1.5), Inches(5.5), Inches(5), LIGHT_GRAY)
prob.shadow.inherit = False

add_textbox(slide, Inches(1), Inches(1.7), Inches(5), Inches(0.5),
    "The Problem", font_size=22, bold=True, color=DARK)

problems = [
    "• 66 days average to form a habit (Lally et al., 2010)",
    "• No easy way to track daily consistency",
    "• Missing visual motivation (streaks, badges)",
    "• Poor analytics — can't identify patterns",
    "• One-size-fits-all — no support for diverse habit types",
    "• Existing apps lack target-count or group features",
]
add_bullet_list(slide, Inches(1), Inches(2.3), Inches(4.8), Inches(3.5), problems, font_size=15)

# Solution box
sol = add_shape(slide, Inches(7), Inches(1.5), Inches(5.5), Inches(5), RGBColor(0xEB, 0xF5, 0xFF))
sol.shadow.inherit = False

add_textbox(slide, Inches(7.2), Inches(1.7), Inches(5), Inches(0.5),
    "Our Solution: HabitForge", font_size=22, bold=True, color=BLUE)

solutions = [
    "✅ 4 habit types: positive, negative, target-count, groups",
    "✅ Daily check-in with optional journal notes",
    "✅ Streak tracking with 6 metrics per habit",
    "✅ Rich analytics: charts, heatmaps, strength scores",
    "✅ Gamification: 6 milestone badges",
    "✅ Mobile-first responsive design with dark mode",
]
add_bullet_list(slide, Inches(7.2), Inches(2.3), Inches(4.8), Inches(3.5), solutions, font_size=15)

# ============================================================
# SLIDE 4: Objectives & Scope
# ============================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide, WHITE)
add_shape(slide, Inches(0), Inches(0), Inches(13.333), Inches(0.08), BLUE)

add_textbox(slide, Inches(0.8), Inches(0.5), Inches(11), Inches(0.8),
    "Objectives & Scope", font_size=32, bold=True, color=DARK)

# Objectives
add_textbox(slide, Inches(0.8), Inches(1.5), Inches(5.5), Inches(0.5),
    "Objectives", font_size=22, bold=True, color=BLUE)

obj_items = [
    "Develop full-stack habit tracking web app",
    "Implement secure authentication (email + Google OAuth)",
    "Design intuitive mobile-first responsive UI",
    "Build analytics dashboard with charts and heatmaps",
    "Incorporate gamification (badges, achievements)",
    "Apply complete SDLC methodology",
]
add_bullet_list(slide, Inches(0.8), Inches(2.1), Inches(5.5), Inches(3.5), obj_items, font_size=15)

# Scope
add_textbox(slide, Inches(7), Inches(1.5), Inches(5.5), Inches(0.5),
    "Scope", font_size=22, bold=True, color=ACCENT_GREEN)

scope_included = [
    "✅ User authentication (email + Google OAuth)",
    "✅ CRUD for 4 habit types",
    "✅ Daily check-in with journal notes",
    "✅ Streak tracking & analytics dashboard",
    "✅ 6 milestone badges",
    "✅ Responsive design + dark/light mode",
    "✅ Deployment on Vercel",
]
add_bullet_list(slide, Inches(7), Inches(2.1), Inches(2.6), Inches(3.5), scope_included, font_size=14)

add_textbox(slide, Inches(9.8), Inches(1.5), Inches(2.5), Inches(0.5),
    "Out of Scope", font_size=18, bold=True, color=ACCENT_ORANGE)

scope_excluded = [
    "❌ Native mobile apps",
    "❌ Social features",
    "❌ Push notifications",
    "❌ Multi-language support",
]
add_bullet_list(slide, Inches(9.8), Inches(2.1), Inches(2.5), Inches(3.5), scope_excluded, font_size=14)

# ============================================================
# SLIDE 5: Technology Stack
# ============================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide, WHITE)
add_shape(slide, Inches(0), Inches(0), Inches(13.333), Inches(0.08), BLUE)

add_textbox(slide, Inches(0.8), Inches(0.5), Inches(11), Inches(0.8),
    "Technology Stack", font_size=32, bold=True, color=DARK)

tech_stack = [
    ("Frontend", "Next.js 14", "React 18, App Router, SSR/CSR hybrid", BLUE),
    ("Styling", "Tailwind CSS v4", "Utility-first, dark mode, responsive", RGBColor(0x06, 0xB6, 0xD4)),
    ("Backend", "Supabase", "BaaS: Auth, Database, Realtime API", ACCENT_GREEN),
    ("Database", "PostgreSQL 15", "Relational, RLS, PostgREST API", ACCENT_PURPLE),
    ("Charts", "Recharts", "React charting library", ACCENT_ORANGE),
    ("Icons", "Lucide React", "Consistent icon system", RGBColor(0xEC, 0x48, 0x99)),
    ("Language", "TypeScript 5", "Type-safe development", RGBColor(0x31, 0x78, 0xC6)),
    ("Testing", "Jest + Playwright", "Unit + E2E testing", RGBColor(0x84, 0xCC, 0x16)),
    ("Deployment", "Vercel", "Serverless, auto-scaling, CDN", RGBColor(0x00, 0x00, 0x00)),
]

y = Inches(1.5)
for category, tech, desc, color in tech_stack:
    # Category badge
    shape = add_shape(slide, Inches(0.8), y, Inches(1.5), Inches(0.45), color)
    add_textbox(slide, Inches(0.8), y + Inches(0.05), Inches(1.5), Inches(0.4),
        category, font_size=12, bold=True, color=WHITE, alignment=PP_ALIGN.CENTER)
    
    # Tech name
    add_textbox(slide, Inches(2.5), y, Inches(2.5), Inches(0.45),
        tech, font_size=18, bold=True, color=TEXT_DARK)
    
    # Description
    add_textbox(slide, Inches(5.2), y, Inches(7), Inches(0.45),
        desc, font_size=14, color=TEXT_GRAY)
    
    y += Inches(0.65)

# ============================================================
# SLIDE 6: System Architecture
# ============================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide, WHITE)
add_shape(slide, Inches(0), Inches(0), Inches(13.333), Inches(0.08), BLUE)

add_textbox(slide, Inches(0.8), Inches(0.5), Inches(11), Inches(0.8),
    "System Architecture", font_size=32, bold=True, color=DARK)

# Client layer
client = add_shape(slide, Inches(0.8), Inches(1.5), Inches(11.7), Inches(1.2), RGBColor(0xEB, 0xF5, 0xFF))
add_textbox(slide, Inches(1), Inches(1.6), Inches(3), Inches(0.4),
    "Client Layer", font_size=16, bold=True, color=BLUE)
add_textbox(slide, Inches(1), Inches(2), Inches(11), Inches(0.5),
    "Browser → React Components (Server + Client) → Tailwind CSS Styling → Recharts Visualization",
    font_size=13, color=TEXT_GRAY)

# Application layer
app = add_shape(slide, Inches(0.8), Inches(2.9), Inches(11.7), Inches(1.5), RGBColor(0xF0, 0xFD, 0xFA))
add_textbox(slide, Inches(1), Inches(3), Inches(3), Inches(0.4),
    "Application Layer", font_size=16, bold=True, color=ACCENT_GREEN)
add_textbox(slide, Inches(1), Inches(3.4), Inches(11), Inches(0.9),
    "Next.js 14 (App Router) on Vercel → Server Components + Client Components → Middleware (Auth) → Supabase Client SDK → Streak Calculator Utility",
    font_size=13, color=TEXT_GRAY)

# Data layer
data = add_shape(slide, Inches(0.8), Inches(4.6), Inches(11.7), Inches(1.2), RGBColor(0xF5, 0xF3, 0xFF))
add_textbox(slide, Inches(1), Inches(4.7), Inches(3), Inches(0.4),
    "Data Layer", font_size=16, bold=True, color=ACCENT_PURPLE)
add_textbox(slide, Inches(1), Inches(5.1), Inches(11), Inches(0.5),
    "Supabase Platform → PostgreSQL 15 (6 tables, RLS) → Auth Service (JWT) → PostgREST API (auto-generated)",
    font_size=13, color=TEXT_GRAY)

# Arrows
for y_pos in [Inches(2.7), Inches(4.4)]:
    add_shape(slide, Inches(6.5), y_pos, Inches(0.3), Inches(0.2), BLUE)

# Key metrics
add_textbox(slide, Inches(0.8), Inches(6.1), Inches(11), Inches(0.4),
    "Key Metrics:", font_size=14, bold=True, color=DARK)

metrics = [
    ("12", "Routes"),
    ("6", "Tables"),
    ("42", "Unit Tests"),
    ("16", "E2E Tests"),
    ("₹0", "Cost"),
    ("100%", "Uptime"),
]

mx = Inches(0.8)
for val, label in metrics:
    add_textbox(slide, mx, Inches(6.5), Inches(1.8), Inches(0.4), val, font_size=20, bold=True, color=BLUE, alignment=PP_ALIGN.CENTER)
    add_textbox(slide, mx, Inches(6.9), Inches(1.8), Inches(0.3), label, font_size=12, color=TEXT_GRAY, alignment=PP_ALIGN.CENTER)
    mx += Inches(1.9)

# ============================================================
# SLIDE 7: Key Features
# ============================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide, WHITE)
add_shape(slide, Inches(0), Inches(0), Inches(13.333), Inches(0.08), BLUE)

add_textbox(slide, Inches(0.8), Inches(0.5), Inches(11), Inches(0.8),
    "Key Features", font_size=32, bold=True, color=DARK)

features = [
    ("✅", "4 Habit Types", "Positive (build), Negative (break), Target Count (numeric), Groups (routines)", BLUE),
    ("📊", "Rich Analytics", "Weekly/monthly charts, category pie chart, 6-month heatmap, habit strength score", ACCENT_GREEN),
    ("🏆", "Gamification", "6 milestone badges: First Step, Week Warrior, Monthly Master, Century Club, Dedicated, Committed", ACCENT_ORANGE),
    ("🌙", "Dark Mode", "Toggle with localStorage persistence + system preference detection", ACCENT_PURPLE),
    ("📱", "Mobile-First", "Fully responsive from 320px to 1920px, touch-friendly 44px targets", RGBColor(0x06, 0xB6, 0xD4)),
    ("🔒", "Secure", "Row Level Security, JWT sessions, HTTPS, XSS/CSRF prevention", RGBColor(0xEF, 0x44, 0x44)),
]

y = Inches(1.5)
for icon, title, desc, color in features:
    # Icon background
    shape = add_shape(slide, Inches(0.8), y, Inches(0.8), Inches(0.8), color)
    add_textbox(slide, Inches(0.8), y + Inches(0.1), Inches(0.8), Inches(0.6),
        icon, font_size=28, color=WHITE, alignment=PP_ALIGN.CENTER)
    
    # Title
    add_textbox(slide, Inches(1.8), y, Inches(3), Inches(0.4),
        title, font_size=18, bold=True, color=TEXT_DARK)
    
    # Description
    add_textbox(slide, Inches(1.8), y + Inches(0.35), Inches(10.5), Inches(0.4),
        desc, font_size=13, color=TEXT_GRAY)
    
    y += Inches(1.0)

# ============================================================
# SLIDE 8: Database Design
# ============================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide, WHITE)
add_shape(slide, Inches(0), Inches(0), Inches(13.333), Inches(0.08), BLUE)

add_textbox(slide, Inches(0.8), Inches(0.5), Inches(11), Inches(0.8),
    "Database Design", font_size=32, bold=True, color=DARK)

# ER diagram simplified
tables = [
    ("habits", "id, user_id, name, habit_type, target_value, frequency, category, color, is_active", "1:N"),
    ("habit_logs", "id, habit_id, user_id, completed_date, count_value, note", "1:N"),
    ("habit_groups", "id, user_id, name, description, color", "M:N"),
    ("habit_group_items", "id, group_id, habit_id, sort_order", ""),
    ("badges", "id, name, description, requirement_type, requirement_value", "1:N"),
    ("user_badges", "id, user_id, badge_id, earned_at", ""),
]

y = Inches(1.5)
for table_name, columns, relation in tables:
    # Table header
    shape = add_shape(slide, Inches(0.8), y, Inches(5.5), Inches(0.4), BLUE)
    add_textbox(slide, Inches(0.9), y + Inches(0.05), Inches(4), Inches(0.35),
        table_name, font_size=14, bold=True, color=WHITE)
    
    if relation:
        add_textbox(slide, Inches(5), y + Inches(0.05), Inches(1), Inches(0.35),
            relation, font_size=11, color=WHITE, alignment=PP_ALIGN.RIGHT)
    
    # Columns
    add_textbox(slide, Inches(0.9), y + Inches(0.4), Inches(5.3), Inches(0.35),
        columns, font_size=10, color=TEXT_GRAY)
    
    y += Inches(0.85)

# Relationships
add_textbox(slide, Inches(7), Inches(1.5), Inches(5.5), Inches(0.5),
    "Relationships", font_size=20, bold=True, color=DARK)

rels = [
    "auth.users ──1:N── habits ──1:N── habit_logs",
    "auth.users ──1:N── habit_groups ──M:N── habits",
    "badges ──1:N── user_badges ◄──1:N── auth.users",
    "",
    "All tables: RLS enabled with ON DELETE CASCADE",
    "Unique constraints: (habit_id, completed_date), (user_id, badge_id)",
]
add_bullet_list(slide, Inches(7), Inches(2.1), Inches(5.5), Inches(3.5), rels, font_size=14)

# ============================================================
# SLIDE 9: Implementation Highlights
# ============================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide, WHITE)
add_shape(slide, Inches(0), Inches(0), Inches(13.333), Inches(0.08), BLUE)

add_textbox(slide, Inches(0.8), Inches(0.5), Inches(11), Inches(0.8),
    "Implementation Highlights", font_size=32, bold=True, color=DARK)

# Left column - Key patterns
add_textbox(slide, Inches(0.8), Inches(1.5), Inches(5.5), Inches(0.5),
    "Key Patterns", font_size=20, bold=True, color=BLUE)

patterns = [
    "• Server Components for data fetching",
    "• Client Components for interactivity",
    "• Middleware for auth session refresh",
    "• Upsert for idempotent check-ins",
    "• Soft delete (is_active flag)",
    "• useMemo for analytics calculation",
    "• Tailwind v4 @custom-variant for dark mode",
]
add_bullet_list(slide, Inches(0.8), Inches(2.1), Inches(5.5), Inches(3.5), patterns, font_size=14)

# Right column - Code snippet
add_textbox(slide, Inches(7), Inches(1.5), Inches(5.5), Inches(0.5),
    "Streak Calculation (TypeScript)", font_size=20, bold=True, color=ACCENT_GREEN)

code = """// Calculate current streak
const uniqueDates = [...new Set(
  logs.map(l => l.completed_date)
)].sort();

let currentStreak = 0;
let checkDate = new Date(startDate);

while (uniqueDates.includes(
  checkDate.toISOString().split('T')[0]
)) {
  currentStreak++;
  checkDate.setDate(
    checkDate.getDate() - 1
  );
}

return { currentStreak, longestStreak,
  completionRate, habitStrength };"""

code_box = add_shape(slide, Inches(7), Inches(2.1), Inches(5.5), Inches(3.5), RGBColor(0x1E, 0x1E, 0x1E))
add_textbox(slide, Inches(7.2), Inches(2.2), Inches(5.1), Inches(3.3),
    code, font_size=11, color=RGBColor(0xD4, 0xD4, 0xD4), font_name="Consolas")

# ============================================================
# SLIDE 10: Testing & Results
# ============================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide, WHITE)
add_shape(slide, Inches(0), Inches(0), Inches(13.333), Inches(0.08), BLUE)

add_textbox(slide, Inches(0.8), Inches(0.5), Inches(11), Inches(0.8),
    "Testing & Results", font_size=32, bold=True, color=DARK)

# Testing pyramid
add_textbox(slide, Inches(0.8), Inches(1.5), Inches(5.5), Inches(0.5),
    "Three-Layer Testing", font_size=20, bold=True, color=BLUE)

# E2E
e2e = add_shape(slide, Inches(3.5), Inches(2.2), Inches(2.5), Inches(0.8), ACCENT_ORANGE)
add_textbox(slide, Inches(3.5), Inches(2.3), Inches(2.5), Inches(0.6),
    "E2E: 16 tests", font_size=16, bold=True, color=WHITE, alignment=PP_ALIGN.CENTER)

# Integration
integ = add_shape(slide, Inches(2.5), Inches(3.2), Inches(4.5), Inches(0.8), ACCENT_GREEN)
add_textbox(slide, Inches(2.5), Inches(3.3), Inches(4.5), Inches(0.6),
    "Component: 12 tests", font_size=16, bold=True, color=WHITE, alignment=PP_ALIGN.CENTER)

# Unit
unit = add_shape(slide, Inches(1.5), Inches(4.2), Inches(6.5), Inches(0.8), BLUE)
add_textbox(slide, Inches(1.5), Inches(4.3), Inches(6.5), Inches(0.6),
    "Unit: 42 tests (Jest)", font_size=16, bold=True, color=WHITE, alignment=PP_ALIGN.CENTER)

# Manual
manual = add_shape(slide, Inches(0.8), Inches(5.2), Inches(8), Inches(0.8), ACCENT_PURPLE)
add_textbox(slide, Inches(0.8), Inches(5.3), Inches(8), Inches(0.6),
    "Manual: 146 checklist items", font_size=16, bold=True, color=WHITE, alignment=PP_ALIGN.CENTER)

# Results summary
add_textbox(slide, Inches(7), Inches(1.5), Inches(5.5), Inches(0.5),
    "Test Results", font_size=20, bold=True, color=DARK)

results = [
    ("42/42", "Unit Tests Passing"),
    ("16/16", "E2E Tests Passing"),
    ("146/146", "Manual Tests Passing"),
    ("14/14", "Visual Screenshots"),
    ("0", "Build Errors"),
    ("12", "Routes Working"),
]

ry = Inches(2.1)
for val, label in results:
    add_textbox(slide, Inches(7), ry, Inches(1.5), Inches(0.4), val, font_size=18, bold=True, color=ACCENT_GREEN)
    add_textbox(slide, Inches(8.5), ry, Inches(3.5), Inches(0.4), label, font_size=14, color=TEXT_DARK)
    ry += Inches(0.55)

# Bugs fixed
add_textbox(slide, Inches(0.8), Inches(6.3), Inches(11), Inches(0.4),
    "7 bugs found and fixed during development (column mismatch, NaN calculations, dark mode persistence, hydration mismatch)",
    font_size=12, color=TEXT_GRAY)

# ============================================================
# SLIDE 11: Live Demo / Screenshots
# ============================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide, WHITE)
add_shape(slide, Inches(0), Inches(0), Inches(13.333), Inches(0.08), BLUE)

add_textbox(slide, Inches(0.8), Inches(0.5), Inches(11), Inches(0.8),
    "Live Demo", font_size=32, bold=True, color=DARK)

# Demo sections
demos = [
    ("🏠", "Dashboard", "5 habits with real-time check-in, progress bars, counter controls"),
    ("📊", "Analytics", "Weekly bar chart, monthly trend, category pie, 6-month heatmap"),
    ("🏆", "Badges", "6 milestone badges with lock/earn status, progress bar"),
    ("📋", "Routines", "Create/manage habit groups with color picker and habit selector"),
    ("🌙", "Dark Mode", "Full dark theme with localStorage persistence"),
    ("📱", "Responsive", "Mobile-first design working from 320px to 1920px"),
]

y = Inches(1.5)
for icon, title, desc in demos:
    shape = add_shape(slide, Inches(0.8), y, Inches(0.7), Inches(0.7), BLUE)
    add_textbox(slide, Inches(0.8), y + Inches(0.08), Inches(0.7), Inches(0.55),
        icon, font_size=24, color=WHITE, alignment=PP_ALIGN.CENTER)
    
    add_textbox(slide, Inches(1.7), y, Inches(2.5), Inches(0.4),
        title, font_size=18, bold=True, color=TEXT_DARK)
    
    add_textbox(slide, Inches(1.7), y + Inches(0.35), Inches(10.5), Inches(0.35),
        desc, font_size=13, color=TEXT_GRAY)
    
    y += Inches(0.85)

# URL
add_textbox(slide, Inches(0.8), Inches(6.5), Inches(11), Inches(0.5),
    "Deployed at: https://habit-forge.vercel.app  |  Source: GitHub Repository",
    font_size=14, color=BLUE, alignment=PP_ALIGN.CENTER)

# ============================================================
# SLIDE 12: Conclusion & Future Scope
# ============================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide, WHITE)
add_shape(slide, Inches(0), Inches(0), Inches(13.333), Inches(0.08), BLUE)

add_textbox(slide, Inches(0.8), Inches(0.5), Inches(11), Inches(0.8),
    "Conclusion & Future Scope", font_size=32, bold=True, color=DARK)

# Conclusion
add_textbox(slide, Inches(0.8), Inches(1.5), Inches(5.5), Inches(0.5),
    "Conclusion", font_size=20, bold=True, color=BLUE)

conclusions = [
    "✅ Complete SDLC execution across 8 Agile sprints",
    "✅ Modern stack: Next.js 14 + Supabase + Tailwind CSS v4",
    "✅ 4 habit types with rich analytics and gamification",
    "✅ 42 unit + 16 E2E + 146 manual tests — all passing",
    "✅ Zero-cost deployment on Vercel + Supabase free tiers",
    "✅ Fully responsive mobile-first design with dark mode",
]
add_bullet_list(slide, Inches(0.8), Inches(2.1), Inches(5.5), Inches(3.5), conclusions, font_size=14)

# Future scope
add_textbox(slide, Inches(7), Inches(1.5), Inches(5.5), Inches(0.5),
    "Future Scope", font_size=20, bold=True, color=ACCENT_GREEN)

future = [
    "🔔 Push notifications for daily reminders",
    "👥 Social features: friends, challenges, leaderboards",
    "📱 React Native mobile app (iOS/Android)",
    "🤖 AI-powered habit insights and recommendations",
    "📅 Google Calendar / Outlook integration",
    "🌐 Multi-language support (Hindi, Punjabi)",
]
add_bullet_list(slide, Inches(7), Inches(2.1), Inches(5.5), Inches(3.5), future, font_size=14)

# Thank you
add_shape(slide, Inches(0), Inches(5.8), Inches(13.333), Inches(1.7), DARK)

add_textbox(slide, Inches(1), Inches(5.9), Inches(11), Inches(0.8),
    "Thank You!", font_size=36, bold=True, color=WHITE, alignment=PP_ALIGN.CENTER)

add_textbox(slide, Inches(1), Inches(6.7), Inches(11), Inches(0.5),
    "Questions? | MCA 4th Semester Major Project | 23ONMCR-753",
    font_size=16, color=TEXT_GRAY, alignment=PP_ALIGN.CENTER)

# Save
output_path = "/Users/openclaw/.hermes/projects/habit-forge/docs/HabitForge_Presentation.pptx"
prs.save(output_path)
print(f"Presentation saved to: {output_path}")
print(f"Slides: {len(prs.slides)}")
