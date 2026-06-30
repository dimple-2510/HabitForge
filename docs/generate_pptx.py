#!/usr/bin/env python3
"""
Generate professional PPT for HabitForge project presentation.
12 slides with consistent design, proper formatting.
"""

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
import copy

# Brand colors
BLUE = RGBColor(0x1E, 0x40, 0xAF)      # Primary blue
BLUE_LIGHT = RGBColor(0x3B, 0x82, 0xF6) # Light blue
BLUE_DARK = RGBColor(0x1E, 0x3A, 0x5F)  # Dark blue
GRAY = RGBColor(0x6B, 0x72, 0x80)       # Gray
GRAY_LIGHT = RGBColor(0xF3, 0xF4, 0xF6) # Light gray
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
BLACK = RGBColor(0x1A, 0x1A, 0x1A)
GREEN = RGBColor(0x10, 0xB9, 0x81)
AMBER = RGBColor(0xF5, 0x9E, 0x0B)
RED = RGBColor(0xEF, 0x44, 0x44)
PURPLE = RGBColor(0x8B, 0x5C, 0xF6)

SLIDE_W = Inches(13.333)
SLIDE_H = Inches(7.5)

def add_bg(slide, color):
    """Set slide background color."""
    background = slide.background
    fill = background.fill
    fill.solid()
    fill.fore_color.rgb = color

def add_shape(slide, left, top, width, height, fill_color, border_color=None):
    """Add a rectangle shape."""
    shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, left, top, width, height)
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill_color
    if border_color:
        shape.line.color.rgb = border_color
        shape.line.width = Pt(1)
    else:
        shape.line.fill.background()
    return shape

def add_rounded_rect(slide, left, top, width, height, fill_color):
    """Add a rounded rectangle."""
    shape = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, left, top, width, height)
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill_color
    shape.line.fill.background()
    return shape

def add_text(slide, left, top, width, height, text, font_size=18, color=BLACK, bold=False, alignment=PP_ALIGN.LEFT, font_name='Segoe UI'):
    """Add a text box."""
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = Pt(font_size)
    p.font.color.rgb = color
    p.font.bold = bold
    p.font.name = font_name
    p.alignment = alignment
    return txBox

def add_bullet_list(slide, left, top, width, height, items, font_size=16, color=BLACK, spacing=Pt(8)):
    """Add a bulleted list."""
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = True
    for i, item in enumerate(items):
        if i == 0:
            p = tf.paragraphs[0]
        else:
            p = tf.add_paragraph()
        p.text = item
        p.font.size = Pt(font_size)
        p.font.color.rgb = color
        p.font.name = 'Segoe UI'
        p.space_after = spacing
        p.level = 0
    return txBox

def add_badge(slide, left, top, text, bg_color=BLUE_LIGHT, text_color=WHITE, font_size=12):
    """Add a small badge/pill shape with text."""
    width = Inches(len(text) * 0.15 + 0.3)
    height = Inches(0.35)
    shape = add_rounded_rect(slide, left, top, width, height, bg_color)
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = False
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = Pt(font_size)
    p.font.color.rgb = text_color
    p.font.name = 'Segoe UI'
    p.font.bold = True
    p.alignment = PP_ALIGN.CENTER
    # Vertically center
    txBox.text_frame.paragraphs[0].space_before = Pt(3)
    return shape

def add_card(slide, left, top, width, height, title, content, icon_text="", bg_color=WHITE, border_color=GRAY_LIGHT):
    """Add a card with title and content."""
    # Card background
    card = add_rounded_rect(slide, left, top, width, height, bg_color)
    card.shadow.inherit = False
    
    # Title
    if icon_text:
        add_text(slide, left + Inches(0.2), top + Inches(0.15), Inches(0.5), Inches(0.4), 
                 icon_text, font_size=20, color=BLUE)
        add_text(slide, left + Inches(0.7), top + Inches(0.1), width - Inches(0.9), Inches(0.5),
                 title, font_size=14, color=BLUE_DARK, bold=True)
    else:
        add_text(slide, left + Inches(0.2), top + Inches(0.1), width - Inches(0.4), Inches(0.5),
                 title, font_size=14, color=BLUE_DARK, bold=True)
    
    # Content
    add_bullet_list(slide, left + Inches(0.2), top + Inches(0.55), width - Inches(0.4), height - Inches(0.7),
                    content, font_size=11, color=GRAY)

def add_header_bar(slide):
    """Add the blue header bar at top of slide."""
    add_shape(slide, 0, 0, SLIDE_W, Inches(0.08), BLUE)

def add_footer(slide, slide_num):
    """Add footer with slide number."""
    add_text(slide, SLIDE_W - Inches(1.5), SLIDE_H - Inches(0.4), Inches(1.3), Inches(0.3),
             f"{slide_num} / 12", font_size=12, color=GRAY, alignment=PP_ALIGN.RIGHT)

def add_section_title(slide, title, subtitle=""):
    """Add section title with decorative element."""
    add_shape(slide, Inches(0.5), Inches(1.2), Inches(0.08), Inches(0.5), BLUE)
    add_text(slide, Inches(0.7), Inches(1.05), Inches(10), Inches(0.7),
             title, font_size=32, color=BLUE_DARK, bold=True)
    if subtitle:
        add_text(slide, Inches(0.7), Inches(1.65), Inches(10), Inches(0.4),
                 subtitle, font_size=16, color=GRAY)

# ==================== SLIDE DEFINITIONS ====================

def slide_title(prs):
    """Slide 1: Title Slide"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])  # Blank
    add_bg(slide, WHITE)
    
    # Blue accent bar at top
    add_shape(slide, 0, 0, SLIDE_W, Inches(0.12), BLUE)
    
    # Decorative blue circle
    circle = slide.shapes.add_shape(MSO_SHAPE.OVAL, Inches(10.5), Inches(1.5), Inches(2.5), Inches(2.5))
    circle.fill.solid()
    circle.fill.fore_color.rgb = BLUE_LIGHT
    circle.line.fill.background()
    
    # Title
    add_text(slide, Inches(1.5), Inches(1.8), Inches(8), Inches(1.2),
             "HabitForge", font_size=54, color=BLUE, bold=True)
    
    add_text(slide, Inches(1.5), Inches(2.9), Inches(8), Inches(0.6),
             "A Full-Stack Habit Tracking Web Application", font_size=24, color=GRAY)
    
    # Tech badges
    badges = ["Next.js 14", "Supabase", "PostgreSQL", "Tailwind CSS v4", "TypeScript", "Vercel"]
    x = Inches(1.5)
    for badge in badges:
        add_badge(slide, x, Inches(3.8), badge, bg_color=BLUE_LIGHT)
        x += Inches(len(badge) * 0.15 + 0.5)
    
    # Meta info
    add_text(slide, Inches(1.5), Inches(4.8), Inches(8), Inches(0.4),
             "MCA 4th Semester Major Project", font_size=18, color=BLUE_DARK, bold=True)
    add_text(slide, Inches(1.5), Inches(5.2), Inches(8), Inches(0.3),
             "Subject Code: 23ONMCR-753  |  12 Credits", font_size=14, color=GRAY)
    add_text(slide, Inches(1.5), Inches(5.5), Inches(8), Inches(0.3),
             "Centre for Distance & Online Education, Chandigarh University", font_size=12, color=GRAY)
    
    # Bottom bar
    add_shape(slide, 0, SLIDE_H - Inches(0.08), SLIDE_W, Inches(0.08), BLUE)

def slide_agenda(prs):
    """Slide 2: Agenda"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_bg(slide, WHITE)
    add_header_bar(slide)
    
    add_text(slide, Inches(0.8), Inches(0.3), Inches(5), Inches(0.6),
             "Agenda", font_size=36, color=BLUE_DARK, bold=True)
    
    items = [
        ("01", "Introduction & Problem Statement", "Why HabitForge?"),
        ("02", "Objectives & Scope", "What we aim to achieve"),
        ("03", "Technology Stack", "Tools and frameworks used"),
        ("04", "System Architecture", "High-level design overview"),
        ("05", "Key Features", "Core functionality walkthrough"),
        ("06", "Database Design", "Schema and relationships"),
        ("07", "Implementation Highlights", "Code snippets and patterns"),
        ("08", "Testing & Results", "Quality assurance approach"),
        ("09", "Live Demo", "Application walkthrough"),
        ("10", "Conclusion & Future Scope", "Summary and roadmap"),
    ]
    
    y = Inches(1.2)
    for num, title, subtitle in items:
        # Number circle
        num_shape = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(1.2), y, Inches(0.6), Inches(0.45))
        num_shape.fill.solid()
        num_shape.fill.fore_color.rgb = BLUE
        num_shape.line.fill.background()
        
        txBox = slide.shapes.add_textbox(Inches(1.2), y + Inches(0.02), Inches(0.6), Inches(0.45))
        tf = txBox.text_frame
        p = tf.paragraphs[0]
        p.text = num
        p.font.size = Pt(14)
        p.font.color.rgb = WHITE
        p.font.bold = True
        p.alignment = PP_ALIGN.CENTER
        
        # Title and subtitle
        add_text(slide, Inches(2.0), y - Inches(0.02), Inches(5), Inches(0.35),
                 title, font_size=16, color=BLUE_DARK, bold=True)
        add_text(slide, Inches(2.0), y + Inches(0.28), Inches(5), Inches(0.25),
                 subtitle, font_size=11, color=GRAY)
        
        y += Inches(0.58)
    
    add_footer(slide, 2)

def slide_intro(prs):
    """Slide 3: Introduction & Problem Statement"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_bg(slide, WHITE)
    add_header_bar(slide)
    add_section_title(slide, "Introduction & Problem Statement")
    
    # Two-column layout
    # Left: Problem
    left_card = add_rounded_rect(slide, Inches(0.5), Inches(2.2), Inches(5.8), Inches(4.5), RGBColor(0xFE, 0xF2, 0xF2))
    add_text(slide, Inches(0.8), Inches(2.4), Inches(5), Inches(0.5),
             "The Problem", font_size=22, color=RED, bold=True)
    
    problems = [
        "66 days average to form a habit (Lally et al., 2010)",
        "No easy way to track daily consistency",
        "Missing visual motivation (streaks, badges)",
        "Poor analytics — can't identify patterns",
        "One-size-fits-all — no support for diverse habit types",
        "Existing apps lack target-count or group features",
    ]
    add_bullet_list(slide, Inches(0.8), Inches(3.0), Inches(5.2), Inches(3.5), problems, font_size=14, color=BLACK)
    
    # Right: Solution
    right_card = add_rounded_rect(slide, Inches(6.8), Inches(2.2), Inches(5.8), Inches(4.5), RGBColor(0xEC, 0xFD, 0xF5))
    add_text(slide, Inches(7.1), Inches(2.4), Inches(5), Inches(0.5),
             "Our Solution: HabitForge", font_size=22, color=GREEN, bold=True)
    
    solutions = [
        "4 habit types: positive, negative, target-count, groups",
        "Daily check-in with optional journal notes",
        "Streak tracking with 6 metrics per habit",
        "Rich analytics: charts, heatmaps, strength scores",
        "Gamification: 6 milestone badges",
        "Mobile-first responsive design with dark mode",
    ]
    add_bullet_list(slide, Inches(7.1), Inches(3.0), Inches(5.2), Inches(3.5), solutions, font_size=14, color=BLACK)
    
    add_footer(slide, 3)

def slide_objectives(prs):
    """Slide 4: Objectives & Scope"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_bg(slide, WHITE)
    add_header_bar(slide)
    add_section_title(slide, "Objectives & Scope")
    
    # Objectives
    add_text(slide, Inches(0.7), Inches(2.2), Inches(5), Inches(0.4),
             "Objectives", font_size=20, color=BLUE, bold=True)
    
    objectives = [
        "Develop full-stack habit tracking web app with 4 habit types",
        "Implement secure authentication (email + Google OAuth)",
        "Design intuitive mobile-first responsive UI",
        "Build analytics dashboard with charts and heatmaps",
        "Incorporate gamification through badge/achievement system",
        "Apply complete SDLC methodology from requirements through deployment",
    ]
    add_bullet_list(slide, Inches(0.7), Inches(2.7), Inches(5.5), Inches(4), objectives, font_size=14, color=BLACK)
    
    # Scope - Included
    add_rounded_rect(slide, Inches(6.8), Inches(2.2), Inches(2.8), Inches(4.5), RGBColor(0xEC, 0xFD, 0xF5))
    add_text(slide, Inches(7.0), Inches(2.4), Inches(2.4), Inches(0.4),
             "Included", font_size=16, color=GREEN, bold=True)
    included = ["Auth (email + OAuth)", "4 habit types", "Daily check-ins", "Streak tracking", "Analytics dashboard", "6 milestone badges", "Dark/light mode", "Vercel deployment"]
    add_bullet_list(slide, Inches(7.0), Inches(2.9), Inches(2.4), Inches(3.5), included, font_size=11, color=BLACK)
    
    # Scope - Excluded
    add_rounded_rect(slide, Inches(10.0), Inches(2.2), Inches(2.8), Inches(4.5), RGBColor(0xFE, 0xF2, 0xF2))
    add_text(slide, Inches(10.2), Inches(2.4), Inches(2.4), Inches(0.4),
             "Excluded", font_size=16, color=RED, bold=True)
    excluded = ["Native mobile apps", "Social features", "Push notifications", "Multi-language support"]
    add_bullet_list(slide, Inches(10.2), Inches(2.9), Inches(2.4), Inches(3.5), excluded, font_size=11, color=BLACK)
    
    add_footer(slide, 4)

def slide_tech_stack(prs):
    """Slide 5: Technology Stack"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_bg(slide, WHITE)
    add_header_bar(slide)
    add_section_title(slide, "Technology Stack")
    
    # Frontend
    add_rounded_rect(slide, Inches(0.5), Inches(2.2), Inches(3.8), Inches(4.5), GRAY_LIGHT)
    add_text(slide, Inches(0.7), Inches(2.4), Inches(3.4), Inches(0.4),
             "Frontend", font_size=18, color=BLUE, bold=True)
    frontend = [
        ("Next.js 14", "App Router, SSR/CSR hybrid"),
        ("React 18", "Server + Client Components"),
        ("Tailwind CSS v4", "Utility-first styling"),
        ("TypeScript 5", "Type-safe development"),
        ("Recharts", "Data visualization"),
    ]
    y = Inches(3.0)
    for name, desc in frontend:
        add_text(slide, Inches(0.7), y, Inches(1.5), Inches(0.3), name, font_size=12, color=BLUE_DARK, bold=True)
        add_text(slide, Inches(2.2), y, Inches(1.8), Inches(0.3), desc, font_size=10, color=GRAY)
        y += Inches(0.55)
    
    # Backend
    add_rounded_rect(slide, Inches(4.5), Inches(2.2), Inches(3.8), Inches(4.5), GRAY_LIGHT)
    add_text(slide, Inches(4.7), Inches(2.4), Inches(3.4), Inches(0.4),
             "Backend / BaaS", font_size=18, color=BLUE, bold=True)
    backend = [
        ("Supabase", "Backend-as-a-Service"),
        ("PostgreSQL", "Relational database"),
        ("PostgREST", "Auto-generated REST API"),
        ("Row Level Security", "Data access control"),
        ("Supabase Auth", "Email + OAuth"),
    ]
    y = Inches(3.0)
    for name, desc in backend:
        add_text(slide, Inches(4.7), y, Inches(1.5), Inches(0.3), name, font_size=12, color=BLUE_DARK, bold=True)
        add_text(slide, Inches(6.2), y, Inches(1.8), Inches(0.3), desc, font_size=10, color=GRAY)
        y += Inches(0.55)
    
    # DevOps
    add_rounded_rect(slide, Inches(8.5), Inches(2.2), Inches(4.2), Inches(4.5), GRAY_LIGHT)
    add_text(slide, Inches(8.7), Inches(2.4), Inches(3.8), Inches(0.4),
             "DevOps & Tools", font_size=18, color=BLUE, bold=True)
    devops = [
        ("Vercel", "Hosting + CI/CD auto-deploy"),
        ("GitHub", "Version control + source"),
        ("VS Code", "IDE with extensions"),
        ("Jest", "Unit testing (42 tests)"),
        ("Playwright", "E2E testing (16 tests)"),
        ("Supabase CLI", "Database management"),
    ]
    y = Inches(3.0)
    for name, desc in devops:
        add_text(slide, Inches(8.7), y, Inches(1.5), Inches(0.3), name, font_size=12, color=BLUE_DARK, bold=True)
        add_text(slide, Inches(10.2), y, Inches(2.2), Inches(0.3), desc, font_size=10, color=GRAY)
        y += Inches(0.55)
    
    add_footer(slide, 5)

def slide_architecture(prs):
    """Slide 6: System Architecture"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_bg(slide, WHITE)
    add_header_bar(slide)
    add_section_title(slide, "System Architecture")
    
    # Three-layer architecture diagram
    layers = [
        ("Client Layer", ["Browser (Chrome, Firefox, Safari)", "React 18 Components", "Tailwind CSS Styling", "Recharts Visualization"], BLUE_LIGHT),
        ("Application Layer", ["Next.js 14 (App Router) on Vercel", "Server + Client Components", "Middleware (Auth)", "Supabase Client SDK"], BLUE),
        ("Data Layer", ["Supabase Platform", "PostgreSQL Database", "Auth Service", "PostgREST API"], BLUE_DARK),
    ]
    
    y = Inches(2.2)
    for title, items, color in layers:
        # Layer box
        add_rounded_rect(slide, Inches(1.5), y, Inches(10.3), Inches(1.3), color)
        add_text(slide, Inches(1.7), y + Inches(0.05), Inches(3), Inches(0.4),
                 title, font_size=16, color=WHITE, bold=True)
        
        # Items in two columns
        for i, item in enumerate(items):
            col = i % 2
            row = i // 2
            add_text(slide, Inches(1.7 + col * 5), y + Inches(0.45 + row * 0.35), Inches(4.8), Inches(0.3),
                     "• " + item, font_size=11, color=WHITE)
        
        y += Inches(1.5)
        
        # Arrow between layers
        if y < Inches(6.5):
            arrow = slide.shapes.add_shape(MSO_SHAPE.DOWN_ARROW, Inches(6.3), y - Inches(0.2), Inches(0.7), Inches(0.2))
            arrow.fill.solid()
            arrow.fill.fore_color.rgb = GRAY
            arrow.line.fill.background()
    
    add_footer(slide, 6)

def slide_features(prs):
    """Slide 7: Key Features"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_bg(slide, WHITE)
    add_header_bar(slide)
    add_section_title(slide, "Key Features")
    
    features = [
        ("Habit Types", ["Positive (build)", "Negative (break)", "Target Count (numeric)", "Groups (routines)"], BLUE),
        ("Daily Check-in", ["One-tap completion", "Optional journal notes", "Counter for target habits", "Date picker for backfill"], GREEN),
        ("Analytics", ["Bar charts (weekly)", "Pie charts (categories)", "Area charts (trends)", "Heatmap calendar", "Habit strength score"], AMBER),
        ("Gamification", ["6 milestone badges", "Streak tracking", "Achievement system", "Progress visualization"], PURPLE),
        ("Recurring Events", ["Custom cycles", "Preset templates", "Calendar integration", "Next occurrence tracking"], RGBColor(0xEC, 0x48, 0x99)),
        ("Mobile-First", ["Responsive design", "Dark/light mode", "Touch-friendly controls", "Bottom tab navigation"], RGBColor(0x06, 0xB6, 0xD4)),
    ]
    
    # 3x2 grid of feature cards
    for i, (title, items, color) in enumerate(features):
        col = i % 3
        row = i // 3
        left = Inches(0.5 + col * 4.1)
        top = Inches(2.2 + row * 2.6)
        
        add_rounded_rect(slide, left, top, Inches(3.8), Inches(2.3), WHITE)
        # Color accent at top
        add_rounded_rect(slide, left, top, Inches(3.8), Inches(0.08), color)
        
        add_text(slide, left + Inches(0.2), top + Inches(0.15), Inches(3.4), Inches(0.4),
                 title, font_size=16, color=color, bold=True)
        add_bullet_list(slide, left + Inches(0.2), top + Inches(0.55), Inches(3.4), Inches(1.6),
                        items, font_size=11, color=BLACK)
    
    add_footer(slide, 7)

def slide_database(prs):
    """Slide 8: Database Design"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_bg(slide, WHITE)
    add_header_bar(slide)
    add_section_title(slide, "Database Design")
    
    # Tables
    tables = [
        ("habits", ["id (PK, UUID)", "user_id (FK)", "name", "habit_type", "target_value", "frequency", "category", "color", "links (JSONB)", "is_active"]),
        ("habit_logs", ["id (PK, UUID)", "habit_id (FK)", "user_id (FK)", "date", "completed", "count_value", "note", "created_at"]),
        ("habit_groups", ["id (PK, UUID)", "user_id (FK)", "name", "description", "color", "is_active"]),
        ("badges", ["id (PK, UUID)", "user_id (FK)", "badge_name", "badge_icon", "earned_at"]),
    ]
    
    # ER Diagram - simplified
    add_text(slide, Inches(0.7), Inches(2.2), Inches(5), Inches(0.4),
             "Entity Relationships", font_size=18, color=BLUE, bold=True)
    
    y = Inches(2.7)
    for table_name, columns in tables:
        add_rounded_rect(slide, Inches(0.7), y, Inches(5.5), Inches(1.0), GRAY_LIGHT)
        add_text(slide, Inches(0.9), y + Inches(0.05), Inches(2), Inches(0.3),
                 table_name, font_size=13, color=BLUE, bold=True)
        cols_text = ", ".join(columns[:5]) + "..."
        add_text(slide, Inches(0.9), y + Inches(0.35), Inches(5), Inches(0.3),
                 cols_text, font_size=9, color=GRAY)
        y += Inches(1.1)
    
    # Relationships
    add_rounded_rect(slide, Inches(7), Inches(2.2), Inches(5.5), Inches(2.5), RGBColor(0xEF, 0xF6, 0xFF))
    add_text(slide, Inches(7.2), Inches(2.4), Inches(5), Inches(0.4),
             "Relationships", font_size=18, color=BLUE, bold=True)
    
    relationships = [
        "users 1:N habits",
        "users 1:N habit_logs",
        "habits 1:N habit_logs",
        "users 1:N habit_groups",
        "habit_groups 1:N habit_group_items",
        "users 1:N badges",
        "users 1:N events",
    ]
    add_bullet_list(slide, Inches(7.2), Inches(2.9), Inches(5), Inches(1.8), relationships, font_size=13, color=BLUE_DARK)
    
    # RLS note
    add_rounded_rect(slide, Inches(7), Inches(5.0), Inches(5.5), Inches(1.5), RGBColor(0xEC, 0xFD, 0xF5))
    add_text(slide, Inches(7.2), Inches(5.1), Inches(5), Inches(0.4),
             "Security: Row Level Security (RLS)", font_size=14, color=GREEN, bold=True)
    add_text(slide, Inches(7.2), Inches(5.5), Inches(5), Inches(0.8),
             "All tables have RLS policies ensuring users can only access their own data. Policies cover SELECT, INSERT, UPDATE, DELETE operations.",
             font_size=11, color=BLACK)
    
    add_footer(slide, 8)

def slide_implementation(prs):
    """Slide 9: Implementation Highlights"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_bg(slide, WHITE)
    add_header_bar(slide)
    add_section_title(slide, "Implementation Highlights")
    
    # Key patterns
    patterns = [
        ("Server Components", "Data fetching happens on the server for better performance and security"),
        ("Client Components", "Interactive elements (forms, buttons) use 'use client' directive"),
        ("Middleware", "Auth session refresh on every request via Next.js middleware"),
        ("Optimistic UI", "Local state updates for instant feedback, no router.refresh() needed"),
        ("Badge Auto-Award", "Supabase RPC function award_badges() runs server-side on check-in"),
    ]
    
    y = Inches(2.2)
    for title, desc in patterns:
        add_rounded_rect(slide, Inches(0.7), y, Inches(11.8), Inches(0.7), GRAY_LIGHT)
        add_text(slide, Inches(0.9), y + Inches(0.05), Inches(3), Inches(0.35),
                 title, font_size=14, color=BLUE, bold=True)
        add_text(slide, Inches(0.9), y + Inches(0.35), Inches(11.4), Inches(0.3),
                 desc, font_size=11, color=GRAY)
        y += Inches(0.8)
    
    # Code snippet
    add_rounded_rect(slide, Inches(0.7), y + Inches(0.2), Inches(11.8), Inches(1.8), RGBColor(0x1F, 0x29, 0x37))
    add_text(slide, Inches(0.9), y + Inches(0.3), Inches(11.4), Inches(0.3),
             "Streak Calculation Engine (lib/streaks.ts)", font_size=12, color=RGBColor(0x10, 0xB9, 0x81), bold=True)
    
    code = """export function calculateStreak(logs: HabitLog[]): StreakResult {
  const sorted = logs.sort((a, b) => b.date.localeCompare(a.date));
  let current = 0, longest = 0, temp = 0;
  for (const log of sorted) {
    if (log.completed) { temp++; longest = Math.max(longest, temp); }
    else { if (current === 0) current = temp; temp = 0; }
  }
  if (current === 0) current = temp;
  return { current, longest, total: logs.filter(l => l.completed).length };
}"""
    add_text(slide, Inches(0.9), y + Inches(0.6), Inches(11.4), Inches(1.3),
             code, font_size=10, color=RGBColor(0xE5, 0xE7, 0xEB), font_name='Consolas')
    
    add_footer(slide, 9)

def slide_testing(prs):
    """Slide 10: Testing & Results"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_bg(slide, WHITE)
    add_header_bar(slide)
    add_section_title(slide, "Testing & Results")
    
    # Testing pyramid - three boxes
    tests = [
        ("Unit Tests", "42 tests (Jest)", "Components, utilities, streak calculation", BLUE_DARK, Inches(4.5)),
        ("Component Tests", "12 tests (Jest)", "HabitCard, forms, interactions", BLUE, Inches(5.5)),
        ("E2E Tests", "16 tests (Playwright)", "Full user flows, auth, CRUD", BLUE_LIGHT, Inches(6.5)),
    ]
    
    for name, count, desc, color, top in tests:
        width = Inches(4.5)
        add_rounded_rect(slide, Inches(4.4), top, width, Inches(1.2), color)
        add_text(slide, Inches(4.6), top + Inches(0.1), Inches(4), Inches(0.4),
                 name, font_size=16, color=WHITE, bold=True)
        add_text(slide, Inches(4.6), top + Inches(0.45), Inches(4), Inches(0.3),
                 count, font_size=13, color=WHITE)
        add_text(slide, Inches(4.6), top + Inches(0.75), Inches(4), Inches(0.3),
                 desc, font_size=10, color=WHITE)
    
    # Results summary
    add_rounded_rect(slide, Inches(0.7), Inches(2.2), Inches(3.3), Inches(4.5), RGBColor(0xEC, 0xFD, 0xF5))
    add_text(slide, Inches(0.9), Inches(2.4), Inches(2.9), Inches(0.4),
             "Results", font_size=18, color=GREEN, bold=True)
    
    results = [
        "70 total tests",
        "100% pass rate",
        ">80% code coverage",
        "16 E2E tests all pass",
        "Manual testing: 146 items",
        "Cross-browser verified",
    ]
    add_bullet_list(slide, Inches(0.9), Inches(2.9), Inches(2.9), Inches(3.5), results, font_size=13, color=BLACK)
    
    # Build output
    add_rounded_rect(slide, Inches(9.3), Inches(2.2), Inches(3.3), Inches(4.5), GRAY_LIGHT)
    add_text(slide, Inches(9.5), Inches(2.4), Inches(2.9), Inches(0.4),
             "Build Output", font_size=18, color=BLUE, bold=True)
    
    build_info = [
        "12 routes compiled",
        "0 TypeScript errors",
        "0 ESLint warnings",
        "First Load JS: 86 KB",
        "Build time: ~15s",
        "Auto-deploy on push",
    ]
    add_bullet_list(slide, Inches(9.5), Inches(2.9), Inches(2.9), Inches(3.5), build_info, font_size=13, color=BLACK)
    
    add_footer(slide, 10)

def slide_demo(prs):
    """Slide 11: Live Demo"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_bg(slide, WHITE)
    add_header_bar(slide)
    add_section_title(slide, "Live Demo", "Application Walkthrough")
    
    # Demo sections
    demos = [
        ("Dashboard", "5 habits with real-time check-in, progress bars, counter controls"),
        ("Analytics", "Pie charts, bar graphs, area trends, heatmap calendar"),
        ("Routines", "Group habits, collective progress, start routine flow"),
        ("Calendar", "Month/week views, color-coded completion dots"),
        ("Events", "Recurring events with custom cycles and presets"),
        ("Mobile", "Bottom tab bar, responsive design, dark mode toggle"),
    ]
    
    # Left column
    for i, (title, desc) in enumerate(demos[:3]):
        y = Inches(2.2 + i * 1.5)
        add_rounded_rect(slide, Inches(0.7), y, Inches(5.5), Inches(1.3), GRAY_LIGHT)
        add_text(slide, Inches(0.9), y + Inches(0.1), Inches(5), Inches(0.35),
                 title, font_size=14, color=BLUE, bold=True)
        add_text(slide, Inches(0.9), y + Inches(0.45), Inches(5), Inches(0.7),
                 desc, font_size=11, color=GRAY)
    
    # Right column
    for i, (title, desc) in enumerate(demos[3:]):
        y = Inches(2.2 + i * 1.5)
        add_rounded_rect(slide, Inches(6.8), y, Inches(5.5), Inches(1.3), GRAY_LIGHT)
        add_text(slide, Inches(7.0), y + Inches(0.1), Inches(5), Inches(0.35),
                 title, font_size=14, color=BLUE, bold=True)
        add_text(slide, Inches(7.0), y + Inches(0.45), Inches(5), Inches(0.7),
                 desc, font_size=11, color=GRAY)
    
    # URL
    add_rounded_rect(slide, Inches(3.5), Inches(6.5), Inches(6.3), Inches(0.6), BLUE)
    add_text(slide, Inches(3.5), Inches(6.55), Inches(6.3), Inches(0.5),
             "Live URL: https://habit-forge-seven.vercel.app", font_size=16, color=WHITE, bold=True, alignment=PP_ALIGN.CENTER)
    
    add_footer(slide, 11)

def slide_conclusion(prs):
    """Slide 12: Conclusion & Future Scope"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_bg(slide, WHITE)
    add_header_bar(slide)
    add_section_title(slide, "Conclusion & Future Scope")
    
    # Conclusion
    add_rounded_rect(slide, Inches(0.7), Inches(2.2), Inches(5.5), Inches(4.5), RGBColor(0xEC, 0xFD, 0xF5))
    add_text(slide, Inches(0.9), Inches(2.4), Inches(5), Inches(0.4),
             "Conclusion", font_size=20, color=GREEN, bold=True)
    
    conclusions = [
        "Complete SDLC execution across 8 Agile sprints",
        "Modern stack: Next.js 14 + Supabase + Tailwind CSS v4",
        "4 habit types with rich analytics and gamification",
        "70 tests with 100% pass rate",
        "Mobile-first responsive design with dark mode",
        "Deployed and live on Vercel",
    ]
    add_bullet_list(slide, Inches(0.9), Inches(2.9), Inches(5), Inches(3.5), conclusions, font_size=14, color=BLACK)
    
    # Future Scope
    add_rounded_rect(slide, Inches(6.8), Inches(2.2), Inches(5.8), Inches(4.5), RGBColor(0xEF, 0xF6, 0xFF))
    add_text(slide, Inches(7.0), Inches(2.4), Inches(5), Inches(0.4),
             "Future Scope", font_size=20, color=BLUE, bold=True)
    
    future = [
        "Native mobile apps (React Native)",
        "Social features (friend groups, challenges)",
        "Push notifications for reminders",
        "AI-powered habit recommendations",
        "Data export (CSV, PDF reports)",
        "Multi-language support (i18n)",
        "Habit sharing and templates",
        "Integration with fitness trackers",
    ]
    add_bullet_list(slide, Inches(7.0), Inches(2.9), Inches(5.4), Inches(3.5), future, font_size=14, color=BLACK)
    
    # Thank you
    add_rounded_rect(slide, Inches(3.5), Inches(6.5), Inches(6.3), Inches(0.6), BLUE)
    add_text(slide, Inches(3.5), Inches(6.55), Inches(6.3), Inches(0.5),
             "Thank You — Questions?", font_size=20, color=WHITE, bold=True, alignment=PP_ALIGN.CENTER)
    
    add_footer(slide, 12)

def main():
    prs = Presentation()
    prs.slide_width = SLIDE_W
    prs.slide_height = SLIDE_H
    
    # Generate all slides
    slide_title(prs)
    slide_agenda(prs)
    slide_intro(prs)
    slide_objectives(prs)
    slide_tech_stack(prs)
    slide_architecture(prs)
    slide_features(prs)
    slide_database(prs)
    slide_implementation(prs)
    slide_testing(prs)
    slide_demo(prs)
    slide_conclusion(prs)
    
    output = '/Users/openclaw/.hermes/projects/habit-forge/docs/HabitForge_Presentation.pptx'
    prs.save(output)
    
    import os
    size = os.path.getsize(output)
    print(f"PPT generated: {output} ({size/1024:.1f} KB)")
    print(f"Slides: {len(prs.slides)}")

if __name__ == '__main__':
    main()
